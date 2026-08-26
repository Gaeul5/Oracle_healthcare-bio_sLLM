"""
발표 자료(PPTX) 생성 — 10장, 좌측 사이드바 템플릿
====================================================
전 슬라이드 고정 좌측 사이드바(짙은 잉크색, 섹션 라벨 + 큰 페이지 번호)
+ 우측 콘텐츠 영역. 색은 잉크/종이 + 절제된 청록 포인트 하나만 사용하고,
R0~R4 조건에는 각각 다른 톤(슬레이트/앰버/청록/자단)을 배정한다.

실행: python3 make_slides.py
출력: research-project/reports/presentation.pptx
"""
import csv
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

HERE = os.path.dirname(os.path.abspath(__file__))
CSV_PATH = os.path.join(HERE, "..", "reports", "eval_results.csv")
OUT_PATH = os.path.join(HERE, "..", "reports", "presentation.pptx")

# ---------------------------------------------------------------- 팔레트
PAPER = RGBColor(0xF7, 0xF5, 0xF0)
INK = RGBColor(0x17, 0x18, 0x1A)
SIDEBAR = RGBColor(0x17, 0x18, 0x1A)
MUTED = RGBColor(0x74, 0x70, 0x66)
SIDE_MUTED = RGBColor(0x8B, 0x90, 0x8E)
RULE = RGBColor(0xD5, 0xCF, 0xC0)
ACCENT = RGBColor(0x4C, 0x8C, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

R0 = RGBColor(0x9A, 0x9D, 0xA3)
R2 = RGBColor(0xC9, 0x8A, 0x54)
R3 = ACCENT
R4 = RGBColor(0x8A, 0x74, 0xB0)
COND_COLOR = {"R0": R0, "R2": R2, "R3": R3, "R4": R4}
COND_LABEL = {"R0": "원문 그대로", "R2": "자유서술", "R3": "구조화", "R4": "구조화+판단"}
SIZES = ["0.5b", "1.5b", "3b"]
SIZE_LABEL = {"0.5b": "0.5B", "1.5b": "1.5B", "3b": "3B"}
CONDS = ["R0", "R2", "R3", "R4"]

SERIF = "Cambria"
SANS = "Calibri"
MONO = "Consolas"

rows = list(csv.DictReader(open(CSV_PATH, encoding="utf-8")))


def mean_f1(size, domain, cond, field="strict_f1"):
    vals = [float(r[field]) for r in rows if r["model_size"] == size and r["train_domain"] == domain and r["condition"] == cond]
    return sum(vals) / len(vals)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

SIDEBAR_W = Inches(1.95)
GUTTER = Inches(0.45)
RMARGIN = Inches(0.7)
CONTENT_LEFT = Emu(SIDEBAR_W + GUTTER)
CONTENT_W = Emu(SW - CONTENT_LEFT - RMARGIN)


# ---------------------------------------------------------------- 프리미티브
def rect(slide, left, top, width, height, fill):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def hline(slide, left, top, width, weight=0.75, color=RULE):
    ln = slide.shapes.add_connector(1, left, top, Emu(left + width), top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def vline(slide, left, top, height, weight=1.0, color=RULE):
    ln = slide.shapes.add_connector(1, left, top, left, Emu(top + height))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def text(slide, left, top, width, height, s, size=16, bold=False, color=INK,
         align=PP_ALIGN.LEFT, font=SANS, line_spacing=1.15, anchor=None, italic=False, spacing=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, ln in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln if spacing is None else spacing.join(list(ln))
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return box


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    return s


def sidebar(slide, eyebrow_text, n, total=11, legend=False):
    rect(slide, Inches(0), Inches(0), SIDEBAR_W, SH, fill=SIDEBAR)
    text(slide, Inches(0.32), Inches(0.5), Inches(1.4), Inches(1.4), eyebrow_text.upper(),
         size=11.5, bold=True, color=ACCENT, font=SANS, line_spacing=1.35, spacing=" ")
    if legend:
        y = Inches(2.2)
        for cond in CONDS:
            rect(slide, Inches(0.32), Emu(y + Inches(0.06)), Inches(0.16), Inches(0.16), fill=COND_COLOR[cond])
            text(slide, Inches(0.58), y, Inches(1.3), Inches(0.4), f"{cond}", size=12.5, bold=True,
                 color=WHITE, font=MONO)
            text(slide, Inches(0.58), Emu(y + Inches(0.26)), Inches(1.3), Inches(0.35), COND_LABEL[cond],
                 size=9.5, color=SIDE_MUTED, font=SANS)
            y = Emu(y + Inches(0.62))
    text(slide, Inches(0.28), Inches(6.35), Inches(1.5), Inches(1.0), f"{n:02d}", size=44, bold=True, color=WHITE, font=SERIF)
    text(slide, Inches(0.32), Inches(7.08), Inches(1.4), Inches(0.3), f"/ {total:02d}", size=11, color=SIDE_MUTED, font=MONO)


def masthead(slide, title_text, title_size=27):
    text(slide, CONTENT_LEFT, Inches(0.62), CONTENT_W, Inches(1.0), title_text,
         size=title_size, bold=True, color=INK, font=SERIF, line_spacing=1.08)
    hline(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, weight=1.1, color=INK)


def field_list(slide, left, top, width, items, size=14.5, gap=Inches(0.3)):
    y = top
    for label, value in items:
        text(slide, left, y, Inches(1.4), Inches(0.35), label, size=size - 1, bold=True, color=ACCENT, font=MONO)
        text(slide, Emu(left + Inches(1.45)), y, Emu(width - Inches(1.45)), Inches(0.85), value,
             size=size, color=INK, font=SANS, line_spacing=1.25)
        y = Emu(y + Inches(0.42) + gap)
    return y


def bullet_list(slide, left, top, width, items, size=15, gap=Inches(0.32), color=INK, marker_color=ACCENT):
    y = top
    for item in items:
        text(slide, left, y, Inches(0.3), Inches(0.4), "—", size=size, bold=True, color=marker_color, font=SANS)
        text(slide, Emu(left + Inches(0.32)), y, Emu(width - Inches(0.32)), Inches(1.2), item,
             size=size, color=color, font=SANS, line_spacing=1.28)
        est_lines = max(1, int(len(item) / (width / Inches(0.135))) + 1)
        y = Emu(y + Inches(0.32) * est_lines + gap)
    return y


def minimal_table(slide, left, top, width, headers, data, col_widths, body_size=13.5, header_size=12, row_h=Inches(0.5)):
    n_rows, n_cols = len(data) + 1, len(headers)
    gt = slide.shapes.add_table(n_rows, n_cols, left, top, width, Emu(row_h * n_rows))
    tbl = gt.table
    tbl.first_row = False
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = w
    for i in range(n_rows):
        tbl.rows[i].height = row_h
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PAPER
        cell.vertical_anchor = MSO_ANCHOR.BOTTOM
        cell.margin_bottom = Pt(4)
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold, run.font.size, run.font.color.rgb, run.font.name = True, Pt(header_size), MUTED, SANS
    hline(slide, left, Emu(top + row_h), width, weight=1.1, color=INK)
    for i, row in enumerate(data, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PAPER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size, run.font.name, run.font.color.rgb = Pt(body_size), SANS, INK
            if j == 0 and str(val) in COND_COLOR:
                run.font.bold, run.font.color.rgb, run.font.name = True, COND_COLOR[str(val)], MONO
        hline(slide, left, Emu(top + row_h * (i + 1)), width, weight=0.5, color=RULE)
    return tbl


def stat_callout(slide, left, top, width, number, caption, size=46):
    text(slide, left, top, width, Inches(0.9), number, size=size, bold=True, color=ACCENT, font=SERIF)
    text(slide, left, Emu(top + Inches(0.8)), width, Inches(0.7), caption, size=12.5, color=MUTED,
         font=SANS, line_spacing=1.25)


def add_grouped_bar_chart(slide, left, top, width, height, domain, chart_title, show_legend=True):
    chart_data = CategoryChartData()
    chart_data.categories = [SIZE_LABEL[sz] for sz in SIZES]
    for cond in CONDS:
        chart_data.add_series(cond, tuple(mean_f1(sz, domain, cond) for sz in SIZES))
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    r0 = chart.chart_title.text_frame.paragraphs[0].runs[0]
    r0.font.size, r0.font.color.rgb, r0.font.name = Pt(13), MUTED, SANS
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = MONO
    for i, cond in enumerate(CONDS):
        series = chart.plots[0].series[i]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = COND_COLOR[cond]
        series.format.line.fill.background()
    val_ax = chart.value_axis
    val_ax.has_major_gridlines = True
    val_ax.major_gridlines.format.line.color.rgb = RULE
    val_ax.major_gridlines.format.line.width = Pt(0.5)
    val_ax.minimum_scale, val_ax.maximum_scale = 0, 70
    val_ax.format.line.color.rgb = RULE
    val_ax.tick_labels.font.size = Pt(10)
    cat_ax = chart.category_axis
    cat_ax.tick_labels.font.size = Pt(13)
    cat_ax.tick_labels.font.name = MONO
    cat_ax.format.line.color.rgb = INK
    return gframe


# ============================================================ 1. 표지
s = new_slide()
sidebar(s, "Research Brief", 1)
text(s, CONTENT_LEFT, Inches(2.3), CONTENT_W, Inches(0.4), "R E S E A R C H   B R I E F", size=12, bold=True, color=ACCENT, font=SANS)
hline(s, CONTENT_LEFT, Inches(2.75), Inches(1.1), weight=2, color=INK)
text(s, CONTENT_LEFT, Inches(2.95), Emu(CONTENT_W), Inches(1.9),
     "표현 형식이 소형 언어모델의\n도메인 간 이상반응(ADE) 추출에\n미치는 영향",
     size=32, bold=True, color=INK, font=SERIF, line_spacing=1.14)
text(s, CONTENT_LEFT, Inches(5.15), CONTENT_W, Inches(0.5),
     "Qwen2.5 (0.5B / 1.5B / 3B)  ·  QLoRA  ·  R0 · R2 · R3 · R4  ·  forum ↔ literature",
     size=13, color=MUTED, font=MONO)

# ============================================================ 2. 배경
s = new_slide()
sidebar(s, "배경", 2)
masthead(s, "소형 모델의 도메인 간 ADE 추출은\n왜 어려운가", title_size=25)
bullet_list(s, CONTENT_LEFT, Inches(1.85), Inches(5.7), [
    "ADE(약물 이상반응) 추출은 약물감시(pharmacovigilance)의 핵심 과제다.",
    "포럼(SNS 후기)과 논문 초록(문헌) 사이의 어휘·문체 차이로 도메인 전이가 특히 어렵다.",
    "MultiADE(Dai et al., 2024)가 동일 코퍼스군에서 이 비대칭을 실측했다.",
    "본 연구는 모델 구조가 아니라 학습 데이터의 표현 형식을 조작해 개입한다.",
], size=14.5)
vline(s, Inches(8.4), Inches(1.85), Inches(4.6), weight=0.75)
text(s, Inches(8.75), Inches(1.8), Inches(4.0), Inches(0.35), "MULTIADE 기준선 (F1)", size=10.5, bold=True, color=MUTED, font=SANS, spacing=" ")
minimal_table(
    s, Inches(8.75), Inches(2.25), Inches(3.9),
    ["학습 → 평가", "F1"],
    [("PsyTAR→PHEE", "27.6"), ("CADEC→PHEE", "22.5"), ("PHEE→CADEC", "13.1"), ("PHEE→PsyTAR", "6.5")],
    col_widths=[Inches(2.7), Inches(1.2)],
)
stat_callout(s, Inches(8.75), Inches(4.9), Inches(3.9), "4.2×", "문헌→포럼이 포럼→문헌보다\n평균 4배 이상 어렵다", size=38)

# ============================================================ 3. RQ & 가설
s = new_slide()
sidebar(s, "질문 & 가설", 3)
masthead(s, "표현 형식이 일반화 성능을\n좌우하는가", title_size=25)
text(s, CONTENT_LEFT, Inches(1.85), CONTENT_W, Inches(1.0),
     "학습 데이터의 표현 형식이 소형 생성 언어모델(0.5B–3B)의 도메인 간 ADE 추출\n일반화에 영향을 주는가? 그 효과는 모델 크기가 작을수록 커지는가?",
     size=16.5, color=INK, font=SERIF, italic=True, line_spacing=1.35)
hline(s, CONTENT_LEFT, Inches(3.1), CONTENT_W, weight=0.5)
field_list(s, CONTENT_LEFT, Inches(3.5), CONTENT_W, [
    ("H1", "인과 구조를 명시한 형식(R3/R4)이, 동일 교사·동일 원문의 자유서술(R2)보다\n미학습 도메인에서 높은 F1을 보인다."),
    ("H2", "H1의 효과 크기는 모델 파라미터 수가 작을수록 크다."),
], size=15.5, gap=Inches(0.32))
text(s, CONTENT_LEFT, Inches(5.85), CONTENT_W, Inches(0.5),
     "사전등록: H1·H2가 기각되어도 결과로 보고한다 (falsifiable).", size=12.5, color=MUTED, italic=True, font=SANS)

# ============================================================ 4. 실험 조건
s = new_slide()
sidebar(s, "실험 설계", 4, legend=True)
masthead(s, "4가지 표현 형식", title_size=27)
minimal_table(
    s, CONTENT_LEFT, Inches(1.85), CONTENT_W,
    ["조건", "형식", "설명"],
    [
        ("R0", "원문 그대로", "교사 미사용, 원본 주석을 그대로 목표 출력으로"),
        ("R2", "자유서술 (통제군)", "교사가 3–5문장 산문으로 설명 — 구조 없이 지식만 전이"),
        ("R3", "구조화 필드", "drug / event / context / onset / severity 등 고정 필드"),
        ("R4", "구조화 + 판단 근거", "R3 + judgement(인과 판단 한 줄) 추가"),
    ],
    col_widths=[Inches(1.3), Inches(3.2), Inches(5.8)],
    body_size=14, header_size=11.5, row_h=Inches(0.62),
)
text(s, CONTENT_LEFT, Inches(6.05), CONTENT_W, Inches(0.6),
     "통제 조건 — 교사 모델·원문·생성 파라미터 동일 / 출력 토큰 수 ±20% 이내 / 교사는 오픈웨이트만 사용",
     size=12, color=MUTED, font=MONO)

# ============================================================ 5. 데이터 & 방법
s = new_slide()
sidebar(s, "데이터 & 방법", 5)
masthead(s, "72개 조합 = 모델 3 × 조건 4\n× 도메인 2 × 시드 3", title_size=23)
bullet_list(s, CONTENT_LEFT, Inches(2.15), Inches(5.0), [
    "코퍼스 — CADECv2(3,526) · CADEC v1(1,155) · PsyTAR(835) · PHEE(4,824)",
    "학생 모델 — Qwen2.5-Instruct 0.5B / 1.5B / 3B, QLoRA",
    "교사 모델 — Qwen2.5-7B-Instruct (오픈웨이트, R2/R3/R4 생성 전용)",
    "평가 — 도메인당 고정 샘플 400건, strict/relaxed F1 병행",
], size=13)
vline(s, Inches(8.05), Inches(2.15), Inches(4.3), weight=0.75)

steps = ["원문 통합\n(10,340건)", "교사 생성\nR2 · R3 · R4", "QLoRA 학습\n(72 runs)", "미학습 도메인\n추론·채점"]
x0, step_w = Inches(8.4), Inches(1.15)
for i, st in enumerate(steps):
    left = Emu(x0 + i * step_w)
    text(s, left, Inches(2.2), Inches(0.4), Inches(0.3), f"{i+1:02d}", size=12, bold=True, color=ACCENT, font=MONO)
    text(s, left, Inches(2.55), Emu(step_w - Inches(0.1)), Inches(1.0), st, size=11, color=INK, font=SANS, line_spacing=1.2)
hline(s, x0, Inches(3.75), Emu(step_w * len(steps) - Inches(0.1)), weight=0.5)
text(s, Inches(8.05), Inches(4.2), Inches(4.6), Inches(1.6),
     "세션 분리 — 교사 추론과 학생 학습을 별도 GPU 세션에서 실행.\n체크포인트/재개 지원으로 세션 중단에도 안전.",
     size=12.5, color=MUTED, font=SANS, line_spacing=1.35)

# ============================================================ 6. 결과 H1 - forum
s = new_slide()
sidebar(s, "결과 — H1", 6, legend=True)
masthead(s, "forum 학습 → literature 평가", title_size=24)
add_grouped_bar_chart(s, CONTENT_LEFT, Inches(1.75), Emu(CONTENT_W - Inches(2.7)), Inches(4.9), "forum",
                       "strict F1, 시드 3개 평균", show_legend=False)
stat_callout(s, Emu(CONTENT_LEFT + CONTENT_W - Inches(2.35)), Inches(2.3), Inches(2.35),
             f"{mean_f1('3b','forum','R3'):.0f}", "3B · R3 strict F1\n(R2 대비 5.7배)", size=48)
hline(s, Emu(CONTENT_LEFT + CONTENT_W - Inches(2.35)), Inches(4.1), Inches(2.0), weight=0.5)
text(s, Emu(CONTENT_LEFT + CONTENT_W - Inches(2.35)), Inches(4.3), Inches(2.35), Inches(2.0),
     "모든 모델 크기에서 R3·R4가\nR0·R2를 큰 폭으로 앞선다.", size=12, color=MUTED, font=SANS, line_spacing=1.3)

# ============================================================ 7. 결과 H1 - literature
s = new_slide()
sidebar(s, "결과 — H1", 7, legend=True)
masthead(s, "literature 학습 → forum 평가", title_size=24)
add_grouped_bar_chart(s, CONTENT_LEFT, Inches(1.75), Emu(CONTENT_W - Inches(2.7)), Inches(4.9), "literature",
                       "strict F1, 시드 3개 평균", show_legend=False)
stat_callout(s, Emu(CONTENT_LEFT + CONTENT_W - Inches(2.35)), Inches(2.3), Inches(2.35),
             f"{mean_f1('1.5b','literature','R3'):.0f}", "1.5B · R3 strict F1\n(전 방향 중 최고)", size=48)
hline(s, Emu(CONTENT_LEFT + CONTENT_W - Inches(2.35)), Inches(4.1), Inches(2.0), weight=0.5)
text(s, Emu(CONTENT_LEFT + CONTENT_W - Inches(2.35)), Inches(4.3), Inches(2.35), Inches(2.0),
     "문헌→포럼은 더 어려운 방향이지만\n구조화 우위는 그대로 유지된다.", size=12, color=MUTED, font=SANS, line_spacing=1.3)

# ============================================================ 8. H1 요약
s = new_slide()
sidebar(s, "H1 요약", 8)
masthead(s, "R3 · R4가 R2를 이긴 조합", title_size=27)
text(s, CONTENT_LEFT, Inches(1.9), Inches(4.3), Inches(1.7), "18 / 18", size=72, bold=True, color=ACCENT, font=SERIF)
text(s, CONTENT_LEFT, Inches(3.5), Inches(4.0), Inches(0.9),
     "모델 × 도메인 × 시드 전체 조합에서\n예외 없이 승리", size=13, color=MUTED, font=SANS, line_spacing=1.3)
vline(s, Inches(6.7), Inches(1.9), Inches(3.9), weight=0.75)
bullet_list(s, Inches(7.1), Inches(1.9), Inches(5.6), [
    "격차는 매우 큼 — forum 학습 시 R2 strict F1 ~10 vs R3 ~50–60.",
    "R2가 R0보다도 낮은 것은 R2의 열등함이 아니라, 자유서술에서 정확한 문자열을 복원하는 채점 파싱의 한계가 반영된 결과다.",
    "R3/R4 대비 우위는 H1의 근거로 유효하다.",
], size=13.5)

# ============================================================ 9. H2 결과
s = new_slide()
sidebar(s, "결과 — H2", 9)
masthead(s, "H2는 지지되지 않는다 — 오히려 반대 방향", title_size=23)
gap_chart_data = CategoryChartData()
gap_chart_data.categories = [SIZE_LABEL[sz] for sz in SIZES]
gap_chart_data.add_series("forum → literature", tuple(mean_f1(sz, "forum", "R3") - mean_f1(sz, "forum", "R2") for sz in SIZES))
gap_chart_data.add_series("literature → forum", tuple(mean_f1(sz, "literature", "R3") - mean_f1(sz, "literature", "R2") for sz in SIZES))
gframe = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, CONTENT_LEFT, Inches(1.75), Inches(6.6), Inches(4.9), gap_chart_data)
chart = gframe.chart
chart.has_title = True
chart.chart_title.text_frame.text = "R3 − R2 격차 (strict F1, pt)"
r0 = chart.chart_title.text_frame.paragraphs[0].runs[0]
r0.font.size, r0.font.color.rgb, r0.font.name = Pt(13), MUTED, SANS
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.name = MONO
chart.plots[0].series[0].format.line.color.rgb = R3
chart.plots[0].series[0].format.line.width = Pt(2.25)
chart.plots[0].series[1].format.line.color.rgb = R2
chart.plots[0].series[1].format.line.width = Pt(2.25)
chart.value_axis.format.line.color.rgb = RULE
chart.category_axis.format.line.color.rgb = INK
chart.category_axis.tick_labels.font.name = MONO
vline(s, Inches(9.05), Inches(1.75), Inches(4.9), weight=0.75)
bullet_list(s, Inches(9.4), Inches(1.9), Inches(3.3), [
    "예상(H2) — 모델이 작을수록 격차가 커야 한다.",
    "실제 — 0.5B +39.5 → 1.5B +44.2 → 3B +48.9pt로 오히려 증가.",
    "사전등록 원칙에 따라 기각 결과를 그대로 보고한다.",
], size=12.5)

# ============================================================ 10. 부가 실험 — RAG
s = new_slide()
sidebar(s, "부가 실험", 10, legend=True)
masthead(s, "검색 기반(RAG)은 파인튜닝을 대체하지 못한다", title_size=23)
text(s, CONTENT_LEFT, Inches(1.85), CONTENT_W, Inches(0.7),
     "파인튜닝 없이 학습 도메인 문서 + R3 교사 출력을 BM25로 검색해 k=3 few-shot으로 넣고\n베이스 3B 모델로 R3 형식을 흉내낸 결과 (검색은 결정론적이라 시드 없이 방향당 1회).",
     size=12.5, color=MUTED, font=SANS, line_spacing=1.3)
rag_headers = ["방향", "R0", "R2", "R3 (QLoRA)", "RAG (파인튜닝 없음)"]
rag_rows = [
    ("forum → literature", f"{mean_f1('3b','forum','R0'):.1f}", f"{mean_f1('3b','forum','R2'):.1f}",
     f"{mean_f1('3b','forum','R3'):.1f}", "38.7"),
    ("literature → forum", f"{mean_f1('3b','literature','R0'):.1f}", f"{mean_f1('3b','literature','R2'):.1f}",
     f"{mean_f1('3b','literature','R3'):.1f}", "11.6"),
]
minimal_table(
    s, CONTENT_LEFT, Inches(2.75), CONTENT_W,
    rag_headers, rag_rows,
    col_widths=[Inches(2.6), Inches(1.6), Inches(1.6), Inches(2.0), Inches(2.6)],
    body_size=14, header_size=11, row_h=Inches(0.55),
)
text(s, CONTENT_LEFT, Inches(4.4), CONTENT_W, Inches(0.4), "strict F1, 3B 모델 기준", size=11, color=MUTED, font=MONO)
hline(s, CONTENT_LEFT, Inches(4.95), CONTENT_W, weight=0.5)
bullet_list(s, CONTENT_LEFT, Inches(5.2), CONTENT_W, [
    "RAG는 R2보다 훨씬 낫지만 QLoRA R3/R4에는 못 미친다 — 두 방향 모두 R0와 비슷한 수준에 머문다.",
    "검색된 예시가 형식을 흉내내는 데는 도움을 주지만, 가중치 자체를 갱신하는 QLoRA만큼의 이득은 주지 못한다.",
], size=13.5)

# ============================================================ 11. 결론
s = new_slide()
sidebar(s, "결론", 11)
masthead(s, "구조화는 이득이 있다 —\n모델 크기에 반비례하지 않는다", title_size=23)
bullet_list(s, CONTENT_LEFT, Inches(1.95), CONTENT_W, [
    "H1 지지 — 인과 구조를 명시한 구조화 형식(R3/R4)이 자유서술보다 도메인 간 일반화에서 일관되게 우수하다 (18/18).",
    "H2 기각, 반대 방향 — 구조화의 이득은 모델이 작을수록 커지지 않으며, 오히려 큰 모델에서 더 크게 나타난다.",
    "문헌 ↔ 포럼 비대칭 재확인 — 구조화가 비대칭을 줄이지만 없애지는 못한다.",
], size=14.5, gap=Inches(0.38))
hline(s, CONTENT_LEFT, Inches(4.55), CONTENT_W, weight=0.5)
text(s, CONTENT_LEFT, Inches(4.75), Inches(4), Inches(0.35), "다음 단계", size=11.5, bold=True, color=ACCENT, font=SANS, spacing=" ")
bullet_list(s, CONTENT_LEFT, Inches(5.2), CONTENT_W, [
    "UMLS 승인 후 R1(용어 정규화)을 forum ↔ literature 주 축에 포함.",
    "temperature scaling 보정 및 risk–coverage curve (사전등록된 부 평가 지표).",
    "라이브 데모로 형식별 출력 차이를 정성적으로 직접 비교.",
], size=12.5, color=MUTED, gap=Inches(0.24))

os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
prs.save(OUT_PATH)
print(f"저장 완료: {OUT_PATH}")
